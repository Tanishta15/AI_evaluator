"""
File Handler utility for AI Idea Evaluator - Phase 1
Handles file format conversions (PPT to PDF) and file management
"""

import os
import shutil
import subprocess
from pathlib import Path
from typing import Dict, Any, Optional
import logging
from datetime import datetime

class FileHandler:
    """Handles file operations including PPT to PDF conversion."""

    def __init__(self, config):
        """Initialize file handler with configuration."""
        self.config = config
        self.logger = logging.getLogger(__name__)
        self.temp_dir = config.temp_dir

        # Supported input formats
        self.ppt_formats = ['.ppt', '.pptx']
        self.pdf_formats = ['.pdf']

        # Conversion statistics
        self.conversion_stats = {
            'total_conversions': 0,
            'successful_conversions': 0,
            'failed_conversions': 0
        }

    def prepare_file_for_extraction(self, file_path: str) -> str:
        """
        Prepare file for extraction by converting PPT to PDF if needed.

        Args:
            file_path: Path to the input file

        Returns:
            Path to the processed file (PDF format)
        """
        input_path = Path(file_path)

        # Validate file exists
        if not input_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        # Check file extension
        file_extension = input_path.suffix.lower()

        if file_extension in self.pdf_formats:
            # Already PDF, no conversion needed
            self.logger.info(f"✅ File is already PDF: {input_path.name}")
            return str(input_path)

        elif file_extension in self.ppt_formats:
            # Need to convert PPT to PDF
            self.logger.info(f"🔄 Converting PPT to PDF: {input_path.name}")
            return self._convert_ppt_to_pdf(str(input_path))

        else:
            # Unsupported format
            raise ValueError(f"Unsupported file format: {file_extension}. Supported formats: {self.config.supported_formats}")

    def _convert_ppt_to_pdf(self, ppt_path: str) -> str:
        """
        Convert PPT/PPTX to PDF using multiple fallback methods.

        Args:
            ppt_path: Path to the PPT file

        Returns:
            Path to the converted PDF file
        """
        input_path = Path(ppt_path)

        # Create output filename
        output_filename = f"{input_path.stem}_converted_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf"
        output_path = self.temp_dir / output_filename

        self.conversion_stats['total_conversions'] += 1

        try:
            # Method 1: Try using LibreOffice (most reliable)
            if self._convert_with_libreoffice(ppt_path, str(output_path)):
                self.conversion_stats['successful_conversions'] += 1
                return str(output_path)

            # Method 2: Try using Python-pptx + reportlab (fallback)
            if self._convert_with_python_libs(ppt_path, str(output_path)):
                self.conversion_stats['successful_conversions'] += 1
                return str(output_path)

            # Method 3: Try using unoconv (another fallback)
            if self._convert_with_unoconv(ppt_path, str(output_path)):
                self.conversion_stats['successful_conversions'] += 1
                return str(output_path)

            # All methods failed
            self.conversion_stats['failed_conversions'] += 1
            raise RuntimeError(f"All conversion methods failed for {ppt_path}")

        except Exception as e:
            self.conversion_stats['failed_conversions'] += 1
            self.logger.error(f"❌ PPT to PDF conversion failed: {str(e)}")
            raise

    def _convert_with_libreoffice(self, input_path: str, output_path: str) -> bool:
        """
        Convert PPT to PDF using LibreOffice headless mode.

        Args:
            input_path: Path to input PPT file
            output_path: Path to output PDF file

        Returns:
            True if conversion successful, False otherwise
        """
        try:
            # Check if LibreOffice is available
            result = subprocess.run(['libreoffice', '--version'], 
                                  capture_output=True, text=True, timeout=10)

            if result.returncode != 0:
                self.logger.warning("LibreOffice not found, skipping this conversion method")
                return False

            # Run LibreOffice conversion
            cmd = [
                'libreoffice',
                '--headless',
                '--convert-to', 'pdf',
                '--outdir', str(self.temp_dir),
                input_path
            ]

            result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)

            if result.returncode == 0:
                # Find the generated PDF file and rename it to our desired output
                input_filename = Path(input_path).stem
                generated_pdf = self.temp_dir / f"{input_filename}.pdf"

                if generated_pdf.exists():
                    shutil.move(str(generated_pdf), output_path)
                    self.logger.info(f"✅ LibreOffice conversion successful: {output_path}")
                    return True
                else:
                    self.logger.warning("LibreOffice conversion completed but PDF not found")
                    return False
            else:
                self.logger.warning(f"LibreOffice conversion failed: {result.stderr}")
                return False

        except (subprocess.TimeoutExpired, subprocess.SubprocessError, Exception) as e:
            self.logger.warning(f"LibreOffice conversion error: {str(e)}")
            return False

    def _convert_with_python_libs(self, input_path: str, output_path: str) -> bool:
        """
        Convert PPT to PDF using Python libraries (python-pptx + reportlab).

        Args:
            input_path: Path to input PPT file
            output_path: Path to output PDF file

        Returns:
            True if conversion successful, False otherwise
        """
        try:
            # This is a simplified conversion using python-pptx and reportlab
            # In a production environment, you might want to use more sophisticated libraries

            from pptx import Presentation
            from reportlab.pdfgen import canvas
            from reportlab.lib.pagesizes import letter, A4
            from reportlab.lib.styles import getSampleStyleSheet
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
            from reportlab.lib.units import inch

            # Load PowerPoint presentation
            prs = Presentation(input_path)

            # Create PDF document
            doc = SimpleDocTemplate(output_path, pagesize=A4)
            styles = getSampleStyleSheet()
            story = []

            # Extract text from each slide
            for i, slide in enumerate(prs.slides):
                # Add slide number
                story.append(Paragraph(f"<b>Slide {i+1}</b>", styles['Heading1']))
                story.append(Spacer(1, 12))

                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        # Clean and add text
                        text = shape.text.strip().replace('\n', '<br/>')
                        story.append(Paragraph(text, styles['Normal']))
                        story.append(Spacer(1, 6))

                # Add page break after each slide (except last)
                if i < len(prs.slides) - 1:
                    story.append(Spacer(1, 24))

            # Build PDF
            doc.build(story)

            self.logger.info(f"✅ Python libs conversion successful: {output_path}")
            return True

        except ImportError:
            self.logger.warning("Required Python libraries not found (python-pptx, reportlab)")
            return False
        except Exception as e:
            self.logger.warning(f"Python libs conversion error: {str(e)}")
            return False

    def _convert_with_unoconv(self, input_path: str, output_path: str) -> bool:
        """
        Convert PPT to PDF using unoconv.

        Args:
            input_path: Path to input PPT file
            output_path: Path to output PDF file

        Returns:
            True if conversion successful, False otherwise
        """
        try:
            # Check if unoconv is available
            result = subprocess.run(['unoconv', '--version'], 
                                  capture_output=True, text=True, timeout=10)

            if result.returncode != 0:
                self.logger.warning("unoconv not found, skipping this conversion method")
                return False

            # Run unoconv conversion
            cmd = [
                'unoconv',
                '-f', 'pdf',
                '-o', output_path,
                input_path
            ]

            result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)

            if result.returncode == 0 and Path(output_path).exists():
                self.logger.info(f"✅ unoconv conversion successful: {output_path}")
                return True
            else:
                self.logger.warning(f"unoconv conversion failed: {result.stderr}")
                return False

        except (subprocess.TimeoutExpired, subprocess.SubprocessError, Exception) as e:
            self.logger.warning(f"unoconv conversion error: {str(e)}")
            return False

    def validate_file(self, file_path: str) -> Dict[str, Any]:
        """
        Validate input file for processing.

        Args:
            file_path: Path to the file to validate

        Returns:
            Dictionary containing validation results
        """
        path = Path(file_path)

        validation_result = {
            'valid': True,
            'errors': [],
            'warnings': [],
            'file_info': {}
        }

        # Check if file exists
        if not path.exists():
            validation_result['valid'] = False
            validation_result['errors'].append(f"File does not exist: {file_path}")
            return validation_result

        # Check file size
        file_size = path.stat().st_size
        validation_result['file_info']['size_bytes'] = file_size
        validation_result['file_info']['size_mb'] = file_size / (1024 * 1024)

        if file_size > self.config.max_file_size:
            validation_result['valid'] = False
            validation_result['errors'].append(
                f"File too large: {file_size / (1024 * 1024):.1f}MB > {self.config.max_file_size / (1024 * 1024):.1f}MB"
            )

        # Check file extension
        file_extension = path.suffix.lower()
        validation_result['file_info']['extension'] = file_extension

        if file_extension not in self.config.supported_formats:
            validation_result['valid'] = False
            validation_result['errors'].append(
                f"Unsupported file format: {file_extension}. Supported: {self.config.supported_formats}"
            )

        # Check if file is readable
        try:
            with open(path, 'rb') as f:
                f.read(1024)  # Try to read first 1KB
        except PermissionError:
            validation_result['valid'] = False
            validation_result['errors'].append("File is not readable (permission denied)")
        except Exception as e:
            validation_result['valid'] = False
            validation_result['errors'].append(f"File read error: {str(e)}")

        return validation_result

    def cleanup_temp_files(self, max_age_hours: int = 24) -> None:
        """
        Clean up temporary files older than specified age.

        Args:
            max_age_hours: Maximum age of temporary files to keep (in hours)
        """
        if not self.temp_dir.exists():
            return

        current_time = datetime.now()
        cleaned_count = 0

        for temp_file in self.temp_dir.glob("*"):
            if temp_file.is_file():
                file_age = current_time - datetime.fromtimestamp(temp_file.stat().st_mtime)

                if file_age.total_seconds() > (max_age_hours * 3600):
                    try:
                        temp_file.unlink()
                        cleaned_count += 1
                        self.logger.debug(f"Cleaned up temp file: {temp_file.name}")
                    except Exception as e:
                        self.logger.warning(f"Failed to clean up {temp_file.name}: {str(e)}")

        if cleaned_count > 0:
            self.logger.info(f"🧹 Cleaned up {cleaned_count} temporary files")

    def get_conversion_stats(self) -> Dict[str, Any]:
        """Get file conversion statistics."""
        return {
            'total_conversions': self.conversion_stats['total_conversions'],
            'successful_conversions': self.conversion_stats['successful_conversions'],
            'failed_conversions': self.conversion_stats['failed_conversions'],
            'success_rate': (
                self.conversion_stats['successful_conversions'] / 
                self.conversion_stats['total_conversions'] 
                if self.conversion_stats['total_conversions'] > 0 else 0
            )
        }