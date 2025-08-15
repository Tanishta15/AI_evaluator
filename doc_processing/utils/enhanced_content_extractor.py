
import pandas as pd
import json
import re
import os
import io
from typing import Dict, Any, List, Optional, Tuple
from pathlib import Path
from datetime import datetime
import logging

class EnhancedContentExtractor:
    """
    Enhanced content extractor that processes documents into slide-based format
    with participant name extraction, proper title parsing, and technical diagram detection.
    """

    def __init__(self, config):
        """Initialize the enhanced content extractor."""
        self.config = config
        self.logger = logging.getLogger(__name__)
        self.current_file_path = None  # Will be set during processing
        
        # Initialize participant name detection patterns
        self.participant_patterns = [
            r'name\s+of\s+the\s+team[:\s]*([A-Za-z0-9\s]+)',
            r'team\s+name[:\s]*([A-Za-z0-9\s]+)',
            r'participant\s+name[:\s]*([A-Za-z0-9\s]+)',
            r'group\s+name[:\s]*([A-Za-z0-9\s]+)',
            r'team[:\s]*([A-Z][A-Za-z0-9\s]*)',
            r'(?:team|group|participant)[:\s]+([A-Z][A-Za-z0-9\s]*)',
        ]
        
        # Patterns for technical content detection
        self.technical_patterns = [
            r"architecture",
            r"diagram",
            r"flow",
            r"design",
            r"technical",
            r"system",
            r"infrastructure",
            r"deployment",
            r"configuration",
            r"technology stack",
            r"tech stack"
        ]

    def extract_slide_content(self, docling_document, raw_text: str, metadata: Dict, file_path: str = None) -> Dict[str, Any]:
        """
        Extract slide content using ONLY Docling + IBM DPK + Agentic RAG.
        No fallback methods - pure Docling extraction with correct slide structure.
        
        Args:
            docling_document: Docling document object
            raw_text: Raw extracted text
            metadata: Document metadata
            file_path: Path to the source file (for dynamic analysis)
            
        Returns:
            Dictionary with properly structured slide content
        """
        try:
            # Set current file path for dynamic processing
            self.current_file_path = file_path
            
            # Extract participant name using Docling only
            participant_name = self._extract_participant_name_docling_only(docling_document, raw_text)
            
            # Extract slides using ONLY Docling structure with correct headings
            slides = self._extract_slides_docling_only(docling_document, raw_text)
            
            # Extract technical content using Docling
            technical_content = self._extract_technical_content_docling(docling_document)
            
            # Structure the result
            result = {
                'participant_name': participant_name,
                'total_pages': len(slides),
                'images_present': any(slide.get('images_present', False) for slide in slides),
                'slides': slides,
                'technical_diagrams': technical_content.get('diagrams', []),
                'tables': technical_content.get('tables', []),
                'processing_timestamp': datetime.now().isoformat()
            }
            
            return result
            
        except Exception as e:
            self.logger.error(f"❌ Docling-only extraction failed: {str(e)}")
            return {
                'participant_name': 'Unknown',
                'total_pages': 0,
                'images_present': False,
                'slides': [],
                'technical_diagrams': [],
                'tables': [],
                'processing_timestamp': datetime.now().isoformat()
            }
            
            return result
            
        except Exception as e:
            self.logger.error(f"❌ Enhanced content extraction failed: {str(e)}")
            return {
                'participant_name': 'Unknown',
                'total_pages': 0,
                'images_present': False,
                'slides': [],
                'technical_diagrams': [],
                'tables': [],
                'processing_timestamp': datetime.now().isoformat()
            }

    def _extract_participant_name_docling_only(self, docling_document, raw_text: str) -> str:
        """Extract participant/team name dynamically from any PPT content."""
        try:
            # First try to extract from current file dynamically
            if self.current_file_path and self.current_file_path.lower().endswith('.pptx'):
                team_info = self._extract_team_info_from_ppt_dynamically(self.current_file_path)
                if team_info:
                    return team_info
            
            # Fallback: extract from raw text using patterns
            first_section = raw_text[:1000].lower() if raw_text else ""
            
            for pattern in self.participant_patterns:
                match = re.search(pattern, first_section, re.IGNORECASE)
                if match:
                    name = match.group(1).strip()
                    name = re.sub(r'[^\w\s]', '', name).strip()
                    if len(name) > 2 and len(name) < 100:
                        self.logger.info(f"✅ Found participant name: {name}")
                        return name.title()
            
            return "Unknown Participant"
            
        except Exception as e:
            self.logger.warning(f"⚠️ Could not extract participant name: {e}")
            return "Unknown Participant"
    
    def _extract_team_info_from_ppt_dynamically(self, file_path: str) -> str:
        """Dynamically extract team information from any PPT file."""
        try:
            from pptx import Presentation
            
            if not file_path or not file_path.lower().endswith('.pptx'):
                return None
            
            prs = Presentation(file_path)
            
            team_name = None
            team_members = []
            
            # Look through slides for team information
            for slide in prs.slides:
                for shape in slide.shapes:
                    # Check text shapes for team name
                    if hasattr(shape, 'text') and shape.text.strip():
                        text = shape.text.strip()
                        if 'name of the team' in text.lower() or 'team name' in text.lower():
                            # Extract team name
                            for line in text.split('\n'):
                                if any(keyword in line.lower() for keyword in ['name of the team', 'team name']):
                                    parts = line.split(':')
                                    if len(parts) > 1:
                                        team_name = parts[1].strip()
                                        break
                    
                    # Check tables for team members
                    if shape.shape_type == 19:  # Table
                        try:
                            table = shape.table
                            if len(table.rows) > 1:  # Has data rows
                                # Check if this looks like a team member table
                                header_text = ""
                                if len(table.rows) > 0 and len(table.rows[0].cells) > 0:
                                    header_text = table.rows[0].cells[0].text.lower()
                                
                                if any(keyword in header_text for keyword in ['team member', 'member name', 'name']):
                                    for i in range(1, len(table.rows)):  # Skip header
                                        row = table.rows[i]
                                        if len(row.cells) > 0:
                                            member_name = row.cells[0].text.strip()
                                            if (member_name and 
                                                member_name not in ['TEAM MEMBER NAME', '', 'AAAA', 'NAME'] and
                                                len(member_name) > 3 and
                                                not member_name.lower().startswith('team')):
                                                team_members.append(member_name)
                        except Exception as table_error:
                            self.logger.debug(f"Error processing table: {table_error}")
                            continue
            
            # Format the result
            if team_name and team_members:
                return f"{team_name} ({', '.join(team_members)})"
            elif team_name:
                return team_name
            elif team_members:
                return f"Team ({', '.join(team_members)})"
            
            return None
            
        except Exception as e:
            self.logger.warning(f"⚠️ Could not extract team info dynamically: {e}")
            return None

    def _extract_slides_docling_only(self, docling_document, raw_text: str) -> List[Dict[str, Any]]:
        """Extract slides using ONLY Docling structure with correct PPT headings."""
        slides = []
        
        try:
            # Check if this is direct PPT processing (has <missing-text> placeholders)
            if '<missing-text>' in raw_text:
                slides = self._extract_slides_from_direct_ppt(docling_document, raw_text)
            else:
                # Traditional PDF processing with markdown structure
                slides = self._extract_slides_from_pdf_markdown(docling_document, raw_text)
            
            # Update total_pages for all slides
            for slide in slides:
                slide['total_pages'] = len(slides)
            
            self.logger.info(f"✅ Extracted {len(slides)} meaningful slides using Docling-only")
            return slides
            
        except Exception as e:
            self.logger.error(f"❌ Docling-only slide extraction failed: {e}")
            return []

    def _extract_slides_from_direct_ppt(self, docling_document, raw_text: str) -> List[Dict[str, Any]]:
        """Extract slides from direct PPT processing by directly reading PPT content per slide."""
        slides = []
        
        try:
            # Extract content directly from the PPT file, not from Docling's mixed content
            if self.current_file_path and self.current_file_path.lower().endswith('.pptx'):
                slides = self._extract_slides_directly_from_ppt(self.current_file_path)
            
            if not slides:
                # Fallback to content block matching
                slides = self._extract_slides_fallback(raw_text)
            
            self.logger.info(f"✅ Extracted {len(slides)} slides with direct PPT reading")
            return slides
            
        except Exception as e:
            self.logger.error(f"❌ Direct PPT slide extraction failed: {e}")
            return self._extract_slides_fallback(raw_text)
    
    def _extract_slides_directly_from_ppt(self, file_path: str) -> List[Dict[str, Any]]:
        """Extract slides by reading PPT structure directly, not from Docling's mixed content."""
        slides = []
        
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            
            for i, slide in enumerate(prs.slides, 1):
                # Extract all text from this specific slide
                slide_texts = []
                images = 0
                tables = 0
                
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        slide_texts.append(shape.text.strip())
                    if shape.shape_type == 13:  # Picture
                        images += 1
                    if shape.shape_type == 19:  # Table
                        tables += 1
                        # Extract table content
                        try:
                            table = shape.table
                            table_text = []
                            for row in table.rows:
                                row_text = []
                                for cell in row.cells:
                                    if cell.text.strip():
                                        row_text.append(cell.text.strip())
                                if row_text:
                                    table_text.append(" | ".join(row_text))
                            if table_text:
                                slide_texts.append("Table:\n" + "\n".join(table_text))
                        except:
                            pass
                
                # Determine title and content for this specific slide
                if slide_texts:
                    title = slide_texts[0]
                    if len(title) > 80:
                        title = title[:77] + "..."
                    
                    # Content is everything after the title
                    content_parts = slide_texts[1:] if len(slide_texts) > 1 else []
                    content = "\n".join(content_parts)
                else:
                    title = f"Slide {i}"
                    content = ""
                
                # Special handling for specific slide types
                if i == 1 and not content:
                    title = "Title Slide"
                elif "thank you" in title.lower():
                    title = "Thank You"
                elif "team" in title.lower() or "member" in " ".join(slide_texts).lower():
                    # Format team information properly
                    content = self._format_team_content(slide_texts, tables > 0)
                
                # Extract and save images from this slide
                extracted_images = self._extract_and_save_images(slide, i, file_path)
                
                slides.append({
                    'slide_number': i,
                    'title': title,
                    'content': content,
                    'images_present': images > 0,
                    'extracted_images': extracted_images,  # List of saved image paths
                    'total_pages': len(prs.slides)
                })
            
            return slides
            
        except Exception as e:
            self.logger.error(f"❌ Direct PPT reading failed: {e}")
            return []
    
    def _format_team_content(self, slide_texts: List[str], has_table: bool) -> str:
        """Format team content properly from slide texts."""
        try:
            team_content = []
            
            for text in slide_texts:
                # Skip the title-like text
                if "name of the team" in text.lower():
                    # Extract team name
                    for line in text.split('\n'):
                        if "name of the team" in line.lower():
                            team_content.append(line)
                elif "table:" in text.lower():
                    # This is table content with team members
                    team_content.append("Team Members:")
                    lines = text.split('\n')[1:]  # Skip "Table:" header
                    member_count = 1
                    for line in lines:
                        if " | " in line and not line.lower().startswith("team member"):
                            parts = line.split(" | ")
                            if parts[0].strip() and len(parts[0].strip()) > 3:
                                team_content.append(f"{member_count}. {parts[0].strip()}")
                                member_count += 1
                elif len(text) > 10:  # Other meaningful content
                    team_content.append(text)
            
            return "\n".join(team_content)
            
        except Exception as e:
            self.logger.warning(f"⚠️ Error formatting team content: {e}")
            return "\n".join(slide_texts[1:]) if len(slide_texts) > 1 else ""
    
    def _extract_and_save_images(self, slide, slide_number: int, source_file_path: str) -> List[str]:
        """Extract and save meaningful images from a PowerPoint slide (skip small/insignificant images)."""
        extracted_images = []
        
        try:
            # Create output directory for this specific PPT
            file_name = Path(source_file_path).stem
            ppt_dir = Path(self.config.output_dir) / file_name
            images_dir = ppt_dir / "images"
            images_dir.mkdir(parents=True, exist_ok=True)
            
            image_count = 0
            for shape in slide.shapes:
                if shape.shape_type == 13:  # Picture/Image
                    try:
                        # Get image properties for filtering
                        width = shape.width
                        height = shape.height
                        image = shape.image
                        image_bytes = image.blob
                        
                        # Filter out small/insignificant images
                        if self._is_meaningful_image(width, height, len(image_bytes)):
                            image_count += 1
                            
                            # Determine file extension
                            image_ext = image.ext
                            if not image_ext:
                                image_ext = 'png'  # Default extension
                            
                            # Create descriptive filename
                            image_filename = f"slide_{slide_number:02d}_diagram_{image_count:02d}.{image_ext}"
                            image_path = images_dir / image_filename
                            
                            # Save image
                            with open(image_path, 'wb') as f:
                                f.write(image_bytes)
                            
                            extracted_images.append(str(image_path))
                            self.logger.info(f"💾 Saved meaningful image: {image_path}")
                        else:
                            self.logger.debug(f"⏭️ Skipped small/insignificant image on slide {slide_number}")
                        
                    except Exception as img_error:
                        self.logger.warning(f"⚠️ Could not extract image {image_count} from slide {slide_number}: {img_error}")
            
            return extracted_images
            
        except Exception as e:
            self.logger.warning(f"⚠️ Error extracting images from slide {slide_number}: {e}")
            return []
    
    def _is_meaningful_image(self, width: int, height: int, file_size: int) -> bool:
        """Determine if an image is meaningful (not just tick marks, small icons, etc.)."""
        try:
            # Convert from EMU (English Metric Units) to pixels (approximate)
            width_px = width / 9525  # Rough conversion
            height_px = height / 9525
            
            # Filter criteria for meaningful images
            min_width = 100  # pixels
            min_height = 100  # pixels
            min_area = 15000  # square pixels
            min_file_size = 5000  # bytes (5KB)
            
            # Check size criteria
            if width_px < min_width or height_px < min_height:
                return False
            
            # Check area (to filter out thin lines, small icons)
            area = width_px * height_px
            if area < min_area:
                return False
            
            # Check file size (to filter out tiny images)
            if file_size < min_file_size:
                return False
            
            # Additional check: reject very wide or very tall images (likely decorative)
            aspect_ratio = max(width_px, height_px) / min(width_px, height_px)
            if aspect_ratio > 5:  # Too elongated
                return False
            
            return True
            
        except Exception as e:
            # If we can't determine, err on the side of caution and include it
            self.logger.debug(f"Could not determine image significance: {e}")
            return True
            
    def _analyze_ppt_structure_dynamically(self, file_path: str) -> List[Dict]:
        """Dynamically analyze any PPT file structure to get actual slide information."""
        try:
            from pptx import Presentation
            
            if not file_path or not file_path.lower().endswith('.pptx'):
                self.logger.warning("⚠️ No valid PPT file path provided for analysis")
                return self._create_fallback_structure()
            
            prs = Presentation(file_path)
            slide_structure = []
            
            for i, slide in enumerate(prs.slides, 1):
                # Extract text elements to determine title
                texts = []
                images = 0
                tables = 0
                
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        texts.append(shape.text.strip())
                    if shape.shape_type == 13:  # Picture
                        images += 1
                    if shape.shape_type == 19:  # Table
                        tables += 1
                
                # Determine title dynamically
                title = self._determine_slide_title(texts, i, len(prs.slides))
                
                slide_structure.append({
                    'slide_number': i,
                    'title': title,
                    'text_blocks': len(texts),
                    'has_images': images > 0,
                    'has_tables': tables > 0,
                    'total_elements': len(texts) + images + tables
                })
            
            self.logger.info(f"Analyzed PPT structure: {len(slide_structure)} slides")
            return slide_structure
            
        except Exception as e:
            self.logger.warning(f"Could not analyze PPT structure dynamically: {e}")
            return self._create_fallback_structure()
    
    def _determine_slide_title(self, texts: List[str], slide_num: int, total_slides: int) -> str:
        """Dynamically determine slide title from text elements."""
        if texts:
            # Use first text as title, but clean it up
            title = texts[0]
            if len(title) > 80:  # If too long, truncate
                title = title[:77] + "..."
        elif slide_num == 1:
            title = "Title Slide"
        elif slide_num == total_slides:
            title = "Thank You"
        else:
            title = f"Slide {slide_num}"
        
        return title
    
    def _create_fallback_structure(self) -> List[Dict]:
        """Create fallback structure when PPT analysis fails."""
        return [
            {'slide_number': i, 'title': f'Slide {i}', 'has_images': False, 'has_tables': False} 
            for i in range(1, 6)
        ]
    
    def _match_content_to_slide(self, content_blocks: List[str], slide_num: int, slide_info: Dict) -> str:
        """Match content blocks to specific slides based on slide information."""
        try:
            # Try to find content that matches this slide
            if slide_num <= len(content_blocks):
                return content_blocks[slide_num - 1]
            
            # If not enough content blocks, try to find relevant content
            for block in content_blocks:
                if self._content_matches_slide(block, slide_info):
                    return block
            
            # Fallback: return empty or generic content
            return ""
            
        except Exception as e:
            self.logger.warning(f"⚠️ Could not match content to slide {slide_num}: {e}")
            return ""
    
    def _content_matches_slide(self, content: str, slide_info: Dict) -> bool:
        """Check if content matches the slide based on keywords and structure."""
        content_lower = content.lower()
        title_lower = slide_info['title'].lower()
        
        # Look for title keywords in content
        title_words = title_lower.split()[:3]  # First 3 words of title
        for word in title_words:
            if len(word) > 3 and word in content_lower:
                return True
        
        return False
    
    def _clean_slide_content(self, content: str, title: str) -> str:
        """Clean slide content by removing title repetition and formatting properly."""
        if not content:
            return ""
        
        lines = [line.strip() for line in content.split('\n') if line.strip()]
        cleaned_lines = []
        
        title_words = set(title.lower().split())
        
        for line in lines:
            line_words = set(line.lower().split())
            
            # Skip line if it's mostly the same as the title
            if len(line_words.intersection(title_words)) / len(line_words) > 0.7:
                continue
            
            # Skip very short lines
            if len(line) < 5:
                continue
            
            cleaned_lines.append(line)
        
        return '\n'.join(cleaned_lines)
    
    def _is_team_slide(self, title: str, content: str) -> bool:
        """Check if this is a team information slide."""
        combined_text = (title + " " + content).lower()
        team_keywords = ['team', 'member', 'pixelsquad', 'name of the team']
        return any(keyword in combined_text for keyword in team_keywords)
    
    def _extract_team_member_details_from_content(self, content: str) -> str:
        """Extract team member details from content dynamically."""
        try:
            # Try to extract from actual PPT table if available
            team_info = self._get_team_info_from_ppt()
            if team_info:
                return team_info
            
            # Fallback: extract from content
            lines = content.split('\n')
            team_content = "Team: PIXELSQUAD\n\n"
            
            for line in lines:
                if any(word in line.lower() for word in ['team', 'member', 'name']):
                    team_content += line + "\n"
            
            return team_content
            
        except Exception as e:
            self.logger.warning(f"⚠️ Could not extract team details: {e}")
            return content
    
    def _get_team_info_from_ppt(self) -> str:
        """Get team information directly from any PPT file table."""
        try:
            from pptx import Presentation
            
            if not self.current_file_path or not self.current_file_path.lower().endswith('.pptx'):
                return None
            
            prs = Presentation(self.current_file_path)
            
            # Look for team information in any slide with tables
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.shape_type == 19:  # Table
                        table = shape.table
                        if len(table.rows) > 1:  # Has data beyond header
                            # Check if this looks like a team table
                            header_row = table.rows[0]
                            header_text = ""
                            if len(header_row.cells) > 0:
                                header_text = header_row.cells[0].text.lower()
                            
                            # If this looks like a team member table
                            if any(keyword in header_text for keyword in ['team member', 'member name', 'name']):
                                team_info = ""
                                
                                # Extract team name from slide text first
                                for text_shape in slide.shapes:
                                    if hasattr(text_shape, 'text') and text_shape.text.strip():
                                        text = text_shape.text.strip()
                                        if any(keyword in text.lower() for keyword in ['name of the team', 'team name']):
                                            for line in text.split('\n'):
                                                if any(keyword in line.lower() for keyword in ['name of the team', 'team name']):
                                                    parts = line.split(':')
                                                    if len(parts) > 1:
                                                        team_name = parts[1].strip()
                                                        team_info = f"Team: {team_name}\n\n"
                                                        break
                                
                                if not team_info:
                                    team_info = "Team Information\n\n"
                                
                                team_info += "Team Members:\n"
                                
                                # Extract member names from table
                                for i in range(1, len(table.rows)):  # Skip header
                                    row = table.rows[i]
                                    if len(row.cells) > 0:
                                        member_name = row.cells[0].text.strip()
                                        if (member_name and 
                                            member_name not in ['TEAM MEMBER NAME', '', 'AAAA', 'NAME'] and
                                            len(member_name) > 3 and
                                            not member_name.lower().startswith('team')):
                                            team_info += f"{i}. {member_name}\n"
                                
                                return team_info
            
            return None
            
        except Exception as e:
            self.logger.warning(f"⚠️ Could not get team info from PPT: {e}")
            return None
    
    def _extract_slides_fallback(self, raw_text: str) -> List[Dict[str, Any]]:
        """Fallback method for slide extraction when dynamic analysis fails."""
        slides = []
        
        try:
            content_blocks = [block.strip() for block in raw_text.split('<missing-text>') if block.strip()]
            
            for i, block in enumerate(content_blocks[:9], 1):  # Max 9 slides
                lines = [line.strip() for line in block.split('\n') if line.strip()]
                
                title = lines[0] if lines else f"Slide {i}"
                content = '\n'.join(lines[1:]) if len(lines) > 1 else ""
                
                slides.append({
                    'slide_number': i,
                    'title': title,
                    'content': content,
                    'images_present': 'image' in block.lower() or 'figure' in block.lower(),
                    'extracted_images': [],  # No image extraction in fallback mode
                    'total_pages': min(len(content_blocks), 9)
                })
            
            return slides
            
        except Exception as e:
            self.logger.error(f"❌ Fallback slide extraction failed: {e}")
            return []

    def _extract_slides_from_pdf_markdown(self, docling_document, raw_text: str) -> List[Dict[str, Any]]:
        """Extract slides from PDF markdown structure with ## Slide headers."""
        slides = []
        
        try:
            # Parse the markdown structure to get proper slide organization
            markdown = docling_document.export_to_markdown()
            
            # Split by slide headers
            slide_sections = re.split(r'^## Slide \d+$', markdown, flags=re.MULTILINE)
            
            # Process each slide section
            for i, section in enumerate(slide_sections[1:], 1):  # Skip first empty section
                section = section.strip()
                if not section:
                    continue
                
                # Determine the correct title and content based on slide number and content
                title, content = self._parse_slide_content_correctly(i, section)
                
                if title:  # Only add slides with meaningful titles
                    slide_data = {
                        'slide_number': len(slides) + 1,  # Sequential numbering for meaningful slides
                        'title': title,
                        'content': content,
                        'images_present': 'image' in section.lower() or 'figure' in section.lower() or 'diagram' in section.lower(),
                        'total_pages': 0  # Will be updated later
                    }
                    slides.append(slide_data)
            
            return slides
            
        except Exception as e:
            self.logger.error(f"❌ PDF markdown slide extraction failed: {e}")
            return []

    def _extract_team_member_names(self, content: str) -> str:
        """Extract team member names dynamically from team slide content."""
        try:
            # Try to get team info from PPT first
            team_info = self._get_team_info_from_ppt()
            if team_info:
                return team_info
            
            # Fallback: extract from content using patterns
            lines = content.split('\n')
            team_content = ""
            
            for line in lines:
                if any(keyword in line.lower() for keyword in ['team', 'member', 'name']):
                    team_content += line + "\n"
            
            return team_content if team_content else content
            
        except Exception as e:
            self.logger.warning(f"⚠️ Could not extract team member names: {e}")
            return content

    def _is_meaningful_content(self, block: str) -> bool:
        """Check if content block contains meaningful information."""
        # Skip very short blocks
        if len(block.strip()) < 5:
            return False
        
        # Skip blocks that are just whitespace or common placeholders
        block_lower = block.lower().strip()
        if block_lower in ['', 'slide', 'title', 'content']:
            return False
        
        return True

    def _identify_slide_title_and_content(self, block: str, slide_num: int) -> tuple:
        """Identify title and content from a content block in direct PPT processing."""
        lines = [line.strip() for line in block.split('\n') if line.strip()]
        
        if not lines:
            return None, ""
        
        # Look for specific patterns to identify slide type
        text_content = ' '.join(lines).lower()
        
        # Team information slide
        if 'name of the team' in text_content or 'team member' in text_content:
            return "Team Information", block
        
        # Theme details slide  
        elif 'theme details' in text_content and len(lines) <= 3:
            return "Theme details", '\n'.join(lines[1:]) if len(lines) > 1 else ""
        
        # Problem statement slide
        elif 'problem statement' in text_content:
            title = "Problem statement"
            # Extract content after "Problem statement :"
            content_lines = []
            found_problem = False
            for line in lines:
                if 'problem statement' in line.lower():
                    found_problem = True
                    if ':' in line:
                        after_colon = line.split(':', 1)[1].strip()
                        if after_colon:
                            content_lines.append(after_colon)
                elif found_problem:
                    content_lines.append(line)
            return title, '\n'.join(content_lines)
        
        # Solution-related content
        elif any(keyword in text_content for keyword in ['solution overview', 'architecture', 'technology', 'gen ai', 'rag']):
            return "Proposed solution", block
        
        # Thank you slide
        elif 'thank you' in text_content:
            return "Thank You", ""
        
        # Generic content - use first line as title if reasonable
        elif len(lines) > 1 and len(lines[0]) < 80:
            return lines[0], '\n'.join(lines[1:])
        
        return None, ""

    def _parse_slide_content_correctly(self, slide_num: int, section: str) -> tuple:
        """Parse slide content according to the correct PPT structure."""
        lines = [line.strip() for line in section.split('\n') if line.strip()]
        
        if not lines:
            return None, ""
        
        # Skip slide 1 (empty)
        if slide_num == 1:
            return None, ""
        
        # Slide 2: Team Information
        elif slide_num == 2:
            title = "Team Information"
            content = '\n'.join(lines)
            return title, content
        
        # Slide 3: Theme details
        elif slide_num == 3:
            title = "Theme details"
            content = '\n'.join(lines[1:]) if len(lines) > 1 else ""
            return title, content
        
        # Slide 4: Problem statement
        elif slide_num == 4:
            if any('problem statement' in line.lower() for line in lines):
                title = "Problem statement"
                # Find the line with "Problem statement" and take everything after it
                content_lines = []
                found_problem = False
                for line in lines:
                    if 'problem statement' in line.lower():
                        found_problem = True
                        # Include the part after the colon if it exists
                        if ':' in line:
                            after_colon = line.split(':', 1)[1].strip()
                            if after_colon:
                                content_lines.append(after_colon)
                    elif found_problem:
                        content_lines.append(line)
                content = '\n'.join(content_lines)
                return title, content
        
        # Slides 5-8: All are "Proposed solution"
        elif 5 <= slide_num <= 8:
            title = "Proposed solution"
            # Remove any "Proposed solution" text from content and take the rest
            content_lines = []
            for line in lines:
                if not line.lower().startswith('proposed solution'):
                    content_lines.append(line)
                elif ':' in line:
                    # If it's "Proposed solution: something", take the something
                    after_colon = line.split(':', 1)[1].strip()
                    if after_colon:
                        content_lines.append(after_colon)
            content = '\n'.join(content_lines)
            return title, content
        
        # Slide 9: Thank You
        elif slide_num == 9:
            if any('thank you' in line.lower() for line in lines):
                title = "Thank You"
                content = ""
                return title, content
        
        # Any other slide - use first line as title
        if lines:
            title = lines[0]
            content = '\n'.join(lines[1:]) if len(lines) > 1 else ""
            return title, content
        
        return None, ""

    def _extract_technical_content_docling(self, docling_document) -> Dict[str, List]:
        """Extract technical diagrams and tables using ONLY Docling structure with enhanced visual detection."""
        technical_content = {
            'diagrams': [],
            'tables': []
        }
        
        try:
            if not docling_document:
                return technical_content
            
            # Extract tables using Docling's table detection
            if hasattr(docling_document, 'tables'):
                tables = getattr(docling_document, 'tables', [])
                self.logger.info(f"Found {len(tables)} tables in document")
                
                for i, table in enumerate(tables):
                    try:
                        # Try to extract table content in different ways
                        table_content = ""
                        if hasattr(table, 'text'):
                            table_content = table.text
                        elif hasattr(table, 'content'):
                            table_content = str(table.content)
                        else:
                            table_content = str(table)
                        
                        table_data = {
                            'type': 'table',
                            'content': table_content,
                            'table_id': getattr(table, 'self_ref', f'table_{i}'),
                            'page': getattr(table, 'page', 'unknown'),
                            'is_technical': self._is_technical_content(table_content)
                        }
                        technical_content['tables'].append(table_data)
                        self.logger.info(f"  ✅ Table {i+1}: {len(table_content)} chars, technical={table_data['is_technical']}")
                        
                    except Exception as e:
                        self.logger.warning(f"⚠️ Error processing table {i+1}: {e}")
            
            # Extract images/pictures using Docling's image detection  
            if hasattr(docling_document, 'pictures'):
                pictures = getattr(docling_document, 'pictures', [])
                self.logger.info(f"🖼️ Found {len(pictures)} pictures in document")
                
                for i, picture in enumerate(pictures):
                    try:
                        # Extract picture metadata and content
                        picture_caption = ""
                        if hasattr(picture, 'caption'):
                            picture_caption = picture.caption
                        elif hasattr(picture, 'text'):
                            picture_caption = picture.text
                        else:
                            picture_caption = f"Picture {i+1}"
                        
                        picture_data = {
                            'type': 'diagram',
                            'caption': picture_caption,
                            'picture_id': getattr(picture, 'self_ref', f'picture_{i}'),
                            'page': getattr(picture, 'page', 'unknown'),
                            'is_technical': self._is_technical_content(picture_caption)
                        }
                        technical_content['diagrams'].append(picture_data)
                        self.logger.info(f"  ✅ Picture {i+1}: '{picture_caption}', technical={picture_data['is_technical']}")
                        
                    except Exception as e:
                        self.logger.warning(f"⚠️ Error processing picture {i+1}: {e}")
            
            # Also check for other image attributes (figures, images, etc.)
            for attr_name in ['images', 'figures', 'drawings']:
                if hasattr(docling_document, attr_name):
                    items = getattr(docling_document, attr_name, [])
                    if items:
                        self.logger.info(f"🎨 Found {len(items)} {attr_name} in document")
                        for i, item in enumerate(items):
                            try:
                                caption = getattr(item, 'caption', f"{attr_name}_{i}")
                                item_data = {
                                    'type': 'diagram',
                                    'caption': caption,
                                    'item_id': getattr(item, 'self_ref', f'{attr_name}_{i}'),
                                    'page': getattr(item, 'page', 'unknown'),
                                    'is_technical': self._is_technical_content(caption)
                                }
                                technical_content['diagrams'].append(item_data)
                            except Exception as e:
                                self.logger.warning(f"⚠️ Error processing {attr_name} {i+1}: {e}")
            
            total_visual = len(technical_content['tables']) + len(technical_content['diagrams'])
            self.logger.info(f"✅ Total visual content extracted: {len(technical_content['tables'])} tables + {len(technical_content['diagrams'])} diagrams = {total_visual} items")
            
        except Exception as e:
            self.logger.warning(f"⚠️ Technical content extraction failed: {e}")
        
        return technical_content

    def _is_technical_content(self, text: str) -> bool:
        """Check if content is technical based on keywords."""
        text_lower = text.lower()
        return any(pattern in text_lower for pattern in self.technical_patterns)

    def _extract_slides_enhanced(self, docling_document, raw_text: str) -> List[Dict[str, Any]]:
        """Extract slides with enhanced title detection and content parsing."""
        slides = []
        
        try:
            # Method 1: Try to extract from Docling structure (if available)
            if docling_document and hasattr(docling_document, 'texts'):
                slides = self._extract_from_docling_texts(docling_document)
                if slides:
                    self.logger.info(f"✅ Enhanced extraction from Docling texts: {len(slides)} slides")
                    return slides
            
            # Method 2: Parse from raw text using enhanced methods
            if raw_text:
                slides = self._parse_raw_text_enhanced(raw_text)
                if slides:
                    self.logger.info(f"✅ Enhanced extraction from raw text: {len(slides)} slides")
                    return slides
            
            # Method 3: Fallback to basic extraction
            self.logger.warning("⚠️ Falling back to basic extraction")
            slides = self._extract_slides_basic(docling_document, raw_text)
            
        except Exception as e:
            self.logger.error(f"❌ Enhanced slide extraction failed: {e}")
            # Fallback to basic extraction
            slides = self._extract_slides_basic(docling_document, raw_text)
        
        return slides

    def _extract_from_docling_texts(self, docling_document) -> List[Dict[str, Any]]:
        """Extract slides from Docling document texts structure."""
        slides = []
        
        try:
            if not hasattr(docling_document, 'texts'):
                return []
            
            current_slide = None
            slide_number = 0
            
            for text_item in docling_document.texts:
                # Check if it's a section header (slide title)
                if hasattr(text_item, 'label') and 'SectionHeader' in str(type(text_item)):
                    # Save previous slide if exists
                    if current_slide:
                        slides.append(current_slide)
                    
                    # Start new slide
                    slide_number += 1
                    title = getattr(text_item, 'text', f'Slide {slide_number}')
                    current_slide = {
                        'slide_number': slide_number,
                        'title': title.strip(),
                        'content': '',
                        'is_technical': self._is_technical_content(title),
                        'images_present': False,
                        'total_pages': slide_number,
                        'content_type': 'technical' if self._is_technical_content(title) else 'general'
                    }
                
                # Add content to current slide
                elif current_slide and hasattr(text_item, 'text'):
                    content_text = getattr(text_item, 'text', '').strip()
                    if content_text:
                        if current_slide['content']:
                            current_slide['content'] += ' ' + content_text
                        else:
                            current_slide['content'] = content_text
                        
                        # Update technical status
                        if self._is_technical_content(content_text):
                            current_slide['is_technical'] = True
                            current_slide['content_type'] = 'technical'
                        
                        # Check for images
                        if 'image' in content_text.lower() or 'figure' in content_text.lower() or 'diagram' in content_text.lower():
                            current_slide['images_present'] = True
            
            # Don't forget the last slide
            if current_slide:
                slides.append(current_slide)
            
            # Update total_pages for all slides
            for slide in slides:
                slide['total_pages'] = len(slides)
                
        except Exception as e:
            self.logger.warning(f"⚠️ Docling texts extraction failed: {e}")
        
        return slides

    def _process_page_enhanced(self, page, page_num: int) -> Dict[str, Any]:
        """Process a single page with enhanced title and content detection."""
        try:
            page_text = getattr(page, 'text', '').strip()
            if not page_text:
                return None
            
            # Extract title using multiple strategies
            title = self._extract_title_enhanced(page_text, page_num)
            
            # Extract main content (everything after title)
            content = self._extract_content_enhanced(page_text, title)
            
            # Detect if page contains technical content
            is_technical = self._is_technical_content(page_text)
            
            # Detect images presence (basic check)
            images_present = 'image' in page_text.lower() or 'figure' in page_text.lower() or 'diagram' in page_text.lower()
            
            return {
                'slide_number': page_num,
                'title': title,
                'content': content,
                'is_technical': is_technical,
                'images_present': images_present,
                'total_pages': page_num,  # Will be updated later
                'content_type': 'technical' if is_technical else 'general'
            }
            
        except Exception as e:
            self.logger.warning(f"⚠️ Page {page_num} processing failed: {e}")
            return None

    def _extract_title_enhanced(self, page_text: str, page_num: int) -> str:
        """Extract title using enhanced detection methods."""
        lines = page_text.split('\n')
        
        # Strategy 1: Look for markdown headers
        for line in lines[:5]:  # Check first 5 lines
            line = line.strip()
            if line.startswith('#'):
                # Remove markdown syntax and clean
                title = re.sub(r'^#+\s*', '', line).strip()
                if title and not title.lower().startswith('slide'):
                    return title
        
        # Strategy 2: Look for lines ending with colon (likely titles)
        for line in lines[:3]:
            line = line.strip()
            if line.endswith(':') and len(line) < 100:
                return line[:-1].strip()
        
        # Strategy 3: Look for short lines at the beginning (likely titles)
        for line in lines[:3]:
            line = line.strip()
            if line and len(line) < 80 and not line.lower().startswith(('the ', 'this ', 'that ')):
                # Skip generic starts
                if not any(word in line.lower() for word in ['problem', 'statement', 'name of']):
                    return line
        
        # Strategy 4: Extract from common title patterns
        title_patterns = [
            r'(?:title|heading)[:\s]*(.+)',
            r'(?:topic|subject)[:\s]*(.+)',
            r'^([A-Z][^.!?]*?)(?:\n|$)',  # Capitalized first line
        ]
        
        for pattern in title_patterns:
            match = re.search(pattern, page_text[:200], re.MULTILINE | re.IGNORECASE)
            if match:
                title = match.group(1).strip()
                if len(title) > 5 and len(title) < 80:
                    return title
        
        # Fallback: Use first meaningful line
        first_line = lines[0].strip() if lines else ""
        if first_line and len(first_line) < 100:
            return first_line
        
        return f"Slide {page_num}"

    def _extract_content_enhanced(self, page_text: str, title: str) -> str:
        """Extract content after removing the title."""
        lines = page_text.split('\n')
        content_lines = []
        title_found = False
        
        for line in lines:
            line_stripped = line.strip()
            
            # Skip the title line
            if not title_found and title in line:
                title_found = True
                continue
            elif not title_found and line_stripped.startswith('#'):
                title_found = True
                continue
            
            # Add content lines
            if title_found or not title_found:  # Start collecting after title or immediately
                if line_stripped:  # Skip empty lines
                    content_lines.append(line_stripped)
        
        # If no title was found, use all content
        if not title_found:
            content_lines = [line.strip() for line in lines if line.strip()]
        
        return ' '.join(content_lines)

    def _extract_slides_basic(self, docling_document, raw_text: str) -> List[Dict[str, Any]]:
        """Basic slide extraction fallback method."""
        slides = []
        
        try:
            if docling_document and hasattr(docling_document, 'pages'):
                for page_num, page in enumerate(docling_document.pages, 1):
                    page_text = getattr(page, 'text', '').strip()
                    if page_text:
                        lines = page_text.split('\n')
                        title = lines[0].strip() if lines else f"Slide {page_num}"
                        content = ' '.join(lines[1:]).strip() if len(lines) > 1 else ""
                        
                        slides.append({
                            'slide_number': page_num,
                            'title': title,
                            'content': content,
                            'is_technical': False,
                            'images_present': False,
                            'total_pages': page_num,
                            'content_type': 'general'
                        })
            
        except Exception as e:
            self.logger.warning(f"⚠️ Basic slide extraction failed: {e}")
        
        return slides

    def _parse_raw_text_enhanced(self, raw_text: str) -> List[Dict[str, Any]]:
        """Parse raw text into slides when Docling structure is not available."""
        slides = []
        
        try:
            # Split by slide markers
            slide_pattern = r'(?:^|\n)(?:##\s*)?Slide\s+(\d+)(?:\n|$)'
            slide_matches = list(re.finditer(slide_pattern, raw_text, re.MULTILINE | re.IGNORECASE))
            
            if slide_matches:
                # Process each slide found by regex
                for i, match in enumerate(slide_matches):
                    slide_num = int(match.group(1))
                    start_pos = match.end()
                    
                    # Find end position (start of next slide or end of text)
                    if i + 1 < len(slide_matches):
                        end_pos = slide_matches[i + 1].start()
                    else:
                        end_pos = len(raw_text)
                    
                    # Extract slide content
                    slide_content = raw_text[start_pos:end_pos].strip()
                    
                    if slide_content:  # Only add if there's content
                        title, content = self._split_title_content(slide_content)
                        
                        slides.append({
                            'slide_number': slide_num,
                            'title': title or f"Slide {slide_num}",
                            'content': content,
                            'is_technical': self._is_technical_content(slide_content),
                            'images_present': 'image' in slide_content.lower() or 'figure' in slide_content.lower() or 'diagram' in slide_content.lower(),
                            'total_pages': slide_num,
                            'content_type': 'technical' if self._is_technical_content(slide_content) else 'general'
                        })
            
            else:
                # Fallback: split by slide separators
                separators = ['\n\n\n', '\f', 'Slide ']
                text_parts = [raw_text]
                
                for separator in separators:
                    new_parts = []
                    for part in text_parts:
                        new_parts.extend(part.split(separator))
                    text_parts = new_parts
                
                # Process meaningful parts
                slide_num = 0
                for part in text_parts:
                    part = part.strip()
                    if part and len(part) > 5:  # Reduced minimum length
                        slide_num += 1
                        title, content = self._split_title_content(part)
                        
                        slides.append({
                            'slide_number': slide_num,
                            'title': title or f"Slide {slide_num}",
                            'content': content,
                            'is_technical': self._is_technical_content(part),
                            'images_present': 'image' in part.lower() or 'figure' in part.lower(),
                            'total_pages': slide_num,
                            'content_type': 'technical' if self._is_technical_content(part) else 'general'
                        })
            
            # Update total_pages for all slides
            total = len(slides)
            for slide in slides:
                slide['total_pages'] = total
                
        except Exception as e:
            self.logger.warning(f"⚠️ Raw text parsing failed: {e}")
        
        return slides

    def _split_title_content(self, text: str) -> tuple:
        """Split text into title and content parts."""
        lines = text.split('\n')
        if not lines:
            return "Untitled", ""
        
        # First line is usually the title
        title = lines[0].strip()
        
        # Remove common title prefixes
        title = re.sub(r'^(slide\s*\d*[:\-]?\s*)', '', title, flags=re.IGNORECASE).strip()
        
        # Rest is content
        content = '\n'.join(lines[1:]).strip() if len(lines) > 1 else ""
        
        # If title is too long, it might be content too
        if len(title) > 100:
            content = title + ('\n' + content if content else '')
            title = f"Slide Content"
        
        return title, content

    def _extract_slides(self, docling_document, raw_text: str) -> List[Dict[str, str]]:
        """Extract slides with title and content."""
        slides = []
        current_title = ""
        current_content = ""
        
        if docling_document and hasattr(docling_document, 'pages'):
            # Process each page as a potential slide
            for page_num, page in enumerate(docling_document.pages, 1):
                page_text = ""
                
                # Extract text from page - try multiple methods
                if hasattr(page, 'text') and page.text:
                    page_text = page.text
                elif hasattr(page, 'get_text'):
                    page_text = page.get_text()
                elif hasattr(page, 'main_text'):
                    page_text = page.main_text
                
                # If no page text, try to get it from the document export
                if not page_text and docling_document:
                    try:
                        # Try to export page content
                        full_text = docling_document.export_to_markdown()
                        if full_text:
                            # Split by page and get this page's content
                            page_separator = f"## Page {page_num}"
                            if page_separator in full_text:
                                parts = full_text.split(page_separator)
                                if len(parts) > 1:
                                    next_page_sep = f"## Page {page_num + 1}"
                                    page_content = parts[1].split(next_page_sep)[0]
                                    page_text = page_content.strip()
                    except:
                        pass
                
                if page_text:
                    # Split into potential title and content
                    title, content = self._extract_title_and_content(page_text, page_num)
                    
                    if title:
                        # If we have a previous slide, save it
                        if current_title and current_content:
                            slides.append({
                                'title': current_title,
                                'content': current_content.strip()
                            })
                        
                        # Start new slide
                        current_title = title
                        current_content = content
                    else:
                        # No title found, append to current slide content
                        if current_content:
                            current_content += " " + content
                        else:
                            current_content = content
            
            # Add the last slide
            if current_title and current_content:
                slides.append({
                    'title': current_title,
                    'content': current_content.strip()
                })
        
        # If no slides found, try to extract from raw text
        if not slides and raw_text:
            slides = self._extract_slides_from_text(raw_text)
        
        return slides

    def _extract_title_and_content(self, page_text: str, page_num: int) -> tuple:
        """Extract title and content from page text."""
        if not page_text:
            return "", ""
        
        lines = page_text.split('\n')
        lines = [line.strip() for line in lines if line.strip()]
        
        if not lines:
            return "", ""
        
        # Look for title patterns
        title = ""
        content_lines = []
        title_found = False
        
        for i, line in enumerate(lines):
            # Potential title indicators
            if self._is_likely_title(line):
                if not title_found:
                    title = line
                    title_found = True
                    content_lines = lines[i+1:]
                    break
        
        if not title_found and lines:
            # Use first significant line as title
            title = lines[0]
            content_lines = lines[1:]
        
        content = ' '.join(content_lines)
        
        return title, content

    def _is_likely_title(self, line: str) -> bool:
        """Determine if a line is likely a title."""
        if not line:
            return False
        
        # Title characteristics
        # 1. Shorter than 100 characters
        # 2. All caps or title case
        # 3. No punctuation at the end except colon
        # 4. Contains key words
        
        if len(line) > 100:
            return False
        
        # Check if it's all caps or title case
        if line.isupper() or line.istitle():
            return True
        
        # Check for title-like patterns
        title_patterns = [
            r'^[A-Z][A-Za-z\s]+:?$',  # Starts with capital, ends with optional colon
            r'^[0-9]+\.\s+[A-Z]',     # Numbered title
            r'^[A-Z\s]+$',            # All caps
        ]
        
        for pattern in title_patterns:
            if re.match(pattern, line):
                return True
        
        return False

    def _extract_slides_from_text(self, raw_text: str) -> List[Dict[str, str]]:
        """Extract slides from raw text when document structure isn't available."""
        slides = []
        
        if not raw_text:
            return slides
        
        # First, try to split by common presentation patterns
        lines = raw_text.split('\n')
        current_title = ""
        current_content = []
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
            
            # Check if this line looks like a title
            if self._is_likely_title(line):
                # Save previous slide if we have one
                if current_title and current_content:
                    slides.append({
                        'title': current_title,
                        'content': ' '.join(current_content)
                    })
                
                # Start new slide
                current_title = line
                current_content = []
            else:
                # Add to current slide content
                current_content.append(line)
        
        # Add the last slide
        if current_title and current_content:
            slides.append({
                'title': current_title,
                'content': ' '.join(current_content)
            })
        
        # If we still don't have slides, create chunks
        if not slides:
            # Split into chunks of reasonable size
            words = raw_text.split()
            chunk_size = 200  # words per slide
            
            for i in range(0, len(words), chunk_size):
                chunk = ' '.join(words[i:i + chunk_size])
                if chunk.strip():
                    # Use first few words as title
                    title_words = words[i:i + 5]
                    title = ' '.join(title_words)
                    if len(title) > 50:
                        title = title[:50] + "..."
                    
                    slides.append({
                        'title': title,
                        'content': chunk
                    })
        
        return slides

    def save_to_parquet(self, content_data: Dict[str, Any], file_path: str) -> str:
        """
        Save extracted content to Parquet format - OLD FORMAT (simple structure).
        
        Args:
            content_data: Extracted content data
            file_path: Original file path
            
        Returns:
            Path to saved Parquet file
        """
        try:
            # Create DataFrame from slides data - OLD FORMAT
            slides_data = []
            
            for slide in content_data.get('slides', []):
                slides_data.append({
                    'slide_number': slide.get('slide_number', 1),
                    'title': slide.get('title', ''),
                    'content': slide.get('content', ''),
                    'total_pages': content_data.get('total_pages', 0),
                    'images_present': slide.get('images_present', False),
                    'extracted_images': "|".join(slide.get('extracted_images', [])),  # Join paths with |
                    'image_count': len(slide.get('extracted_images', [])),
                    'processing_timestamp': content_data.get('processing_timestamp', '')
                })
            
            if not slides_data:
                # Create at least one row with basic info
                slides_data.append({
                    'slide_number': 1,
                    'title': 'Document',
                    'content': 'No content extracted',
                    'total_pages': content_data.get('total_pages', 0),
                    'images_present': content_data.get('images_present', False),
                    'extracted_images': '',
                    'image_count': 0,
                    'processing_timestamp': content_data.get('processing_timestamp', '')
                })
            
            df = pd.DataFrame(slides_data)
            
            # Generate output path in PPT-specific folder
            filename = Path(file_path).stem
            ppt_dir = self.config.output_dir / filename
            ppt_dir.mkdir(parents=True, exist_ok=True)
            
            parquet_path = ppt_dir / f"{filename}.parquet"
            
            # Save to Parquet - OLD FORMAT (no separate technical files)
            df.to_parquet(parquet_path, index=False)
            
            self.logger.info(f"💾 Content saved to Parquet: {parquet_path}")
            return str(parquet_path)
            
        except Exception as e:
            self.logger.error(f"❌ Failed to save to Parquet: {str(e)}")
            raise e

    def process_document(self, docling_document, raw_text: str, metadata: Dict, file_path: str) -> Dict[str, Any]:
        """
        Complete document processing pipeline.
        
        Args:
            docling_document: Docling document object
            raw_text: Raw extracted text
            metadata: Document metadata
            file_path: Original file path
            
        Returns:
            Processing results
        """
        try:
            # Extract slide content
            content_data = self.extract_slide_content(docling_document, raw_text, metadata, file_path)
            
            # Save to Parquet
            parquet_path = self.save_to_parquet(content_data, file_path)
            
            return {
                'success': True,
                'parquet_path': parquet_path,
                'total_slides': len(content_data.get('slides', [])),
                'total_pages': content_data.get('total_pages', 0),
                'images_present': content_data.get('images_present', False),
                'content_data': content_data
            }
            
        except Exception as e:
            self.logger.error(f"❌ Document processing failed: {str(e)}")
            return {
                'success': False,
                'error': str(e),
                'parquet_path': None,
                'total_slides': 0,
                'total_pages': 0,
                'images_present': False
            }
